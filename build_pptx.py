"""
Build round_2/presentation.pptx from DECK.md content.
Run with: /Users/valletivarish/Desktop/flipkart_ml/.venv/bin/python round_2/build_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
BTP_DARK   = RGBColor(0x1A, 0x23, 0x3A)   # deep navy  — title bars
BTP_ACCENT = RGBColor(0xFF, 0x6B, 0x00)   # Flipkart orange — accents
BTP_LIGHT  = RGBColor(0xF5, 0xF7, 0xFA)   # near-white slide bg
BTP_TEXT   = RGBColor(0x1A, 0x23, 0x3A)   # body text
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY_MID   = RGBColor(0x55, 0x65, 0x7A)   # sub-text / secondary

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

BASE = os.path.dirname(os.path.abspath(__file__))
FIG  = os.path.join(BASE, "eda")

# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------

def add_solid_fill(shape, colour: RGBColor):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = colour


def set_text(tf, text, size=18, bold=False, colour=None, align=PP_ALIGN.LEFT):
    tf.text = ""
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if colour:
        run.font.color.rgb = colour


def add_title_bar(slide, title_text, subtitle_text=""):
    """Dark bar across the top with white title text."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), SLIDE_W, Inches(1.35)
    )
    add_solid_fill(bar, BTP_DARK)
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left  = Inches(0.35)
    tf.margin_top   = Inches(0.15)
    tf.margin_bottom= Inches(0.05)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size  = Pt(28)
    run.font.bold  = True
    run.font.color.rgb = WHITE

    if subtitle_text:
        from pptx.util import Pt as _Pt
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle_text
        r2.font.size  = _Pt(14)
        r2.font.bold  = False
        r2.font.color.rgb = RGBColor(0xB8, 0xC5, 0xD6)


def add_accent_line(slide):
    """Thin orange horizontal rule under the title bar."""
    line = slide.shapes.add_shape(
        1,
        Inches(0), Inches(1.35), SLIDE_W, Inches(0.04)
    )
    add_solid_fill(line, BTP_ACCENT)
    line.line.fill.background()


def add_slide_bg(slide):
    """Light background rectangle behind everything."""
    bg = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0), SLIDE_W, SLIDE_H
    )
    add_solid_fill(bg, BTP_LIGHT)
    bg.line.fill.background()
    # send to back
    sp = bg._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def bullet_box(slide, left, top, width, height, items,
                size=17, title=None, title_size=18):
    """
    Add a text box with an optional bold title and bullet list.
    items: list of strings
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(0.1)
    tf.margin_top    = Inches(0.08)
    tf.margin_right  = Inches(0.1)
    tf.margin_bottom = Inches(0.08)

    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size  = Pt(title_size)
        run.font.bold  = True
        run.font.color.rgb = BTP_DARK

    for item in items:
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        # detect whether item starts with a sub-marker
        if item.startswith("    "):
            run.text = "      " + item.strip()
            run.font.size  = Pt(size - 1)
            run.font.color.rgb = GREY_MID
        else:
            run.text = item
            run.font.size  = Pt(size)
            run.font.color.rgb = BTP_TEXT


def add_image(slide, img_path, left, top, width, height):
    """Add an image; skip gracefully if file missing."""
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, left, top, width, height)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_01_title(prs):
    """Cover slide."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Full navy background
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    add_solid_fill(bg, BTP_DARK)
    bg.line.fill.background()

    # Orange accent strip on left
    strip = slide.shapes.add_shape(1, 0, 0, Inches(0.22), SLIDE_H)
    add_solid_fill(strip, BTP_ACCENT)
    strip.line.fill.background()

    # Main title
    tx1 = slide.shapes.add_textbox(Inches(0.55), Inches(1.5), Inches(8.5), Inches(1.5))
    tf1 = tx1.text_frame
    p = tf1.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "Event Intelligence for BTP"
    run.font.size  = Pt(46)
    run.font.bold  = True
    run.font.color.rgb = WHITE

    # Subtitle
    tx2 = slide.shapes.add_textbox(Inches(0.55), Inches(3.1), Inches(9), Inches(0.7))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "From Reactive Response to Proactive Deployment"
    run2.font.size  = Pt(22)
    run2.font.italic = True
    run2.font.color.rgb = RGBColor(0xB8, 0xC5, 0xD6)

    # Key details box
    details = [
        "Dataset:  8,173 real ASTraM events  ·  Bengaluru  ·  Nov 2023 – Apr 2024",
        "Three pillars:  Forecast Impact  →  Recommend Resources  →  Keep Learning",
        "Live prototype:  Streamlit dashboard  +  real-time Event Simulator",
        "Every number traces back to BTP's own ASTraM records",
    ]
    tx3 = slide.shapes.add_textbox(Inches(0.55), Inches(4.05), Inches(11), Inches(2.5))
    tf3 = tx3.text_frame
    tf3.word_wrap = True
    first = True
    for d in details:
        p = tf3.paragraphs[0] if first else tf3.add_paragraph()
        first = False
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = "▸  " + d
        run.font.size  = Pt(17)
        run.font.color.rgb = RGBColor(0xD0, 0xDC, 0xEA)

    # Badge bottom-right
    tx4 = slide.shapes.add_textbox(Inches(9.2), Inches(6.5), Inches(4), Inches(0.7))
    tf4 = tx4.text_frame
    p4 = tf4.paragraphs[0]
    p4.alignment = PP_ALIGN.RIGHT
    r4 = p4.add_run()
    r4.text = "Gridlock Hackathon 2.0  ·  PS2  ·  Bengaluru Traffic Police"
    r4.font.size  = Pt(11)
    r4.font.color.rgb = RGBColor(0x6B, 0x7F, 0x99)

    return slide


def slide_02_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_title_bar(slide, "The Problem: Three Gaps BTP Faces Today",
                  "BTP's own brief names the gaps — we built exactly one pillar per gap")
    add_accent_line(slide)

    # Gap table as styled text boxes
    gaps = [
        ("GAP 1 — Impact Not Quantified",
         "Officers arrive with no estimate of how long an event will last\n"
         "or whether the road will close. Every incident looks equally urgent."),
        ("GAP 2 — Experience-Driven Deployment",
         "Manpower decisions depend on individual judgment — inconsistent\n"
         "across shifts and zones. No principled resource model exists."),
        ("GAP 3 — No Post-Event Learning",
         "Each event is forgotten once cleared. Patterns repeat. The system\n"
         "never improves. The event mix keeps shifting undetected."),
    ]

    tops  = [Inches(1.55), Inches(3.15), Inches(4.75)]
    for i, (heading, body) in enumerate(gaps):
        # coloured label
        lbl = slide.shapes.add_shape(1,
            Inches(0.4), tops[i], Inches(0.08), Inches(1.3))
        add_solid_fill(lbl, BTP_ACCENT)
        lbl.line.fill.background()

        tx = slide.shapes.add_textbox(
            Inches(0.65), tops[i], Inches(12.3), Inches(1.35))
        tf = tx.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = heading
        r1.font.size  = Pt(18)
        r1.font.bold  = True
        r1.font.color.rgb = BTP_DARK
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = body
        r2.font.size  = Pt(16)
        r2.font.color.rgb = GREY_MID

    # Stats strip at bottom
    stats_box = slide.shapes.add_textbox(Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.85))
    tf_s = stats_box.text_frame
    tf_s.word_wrap = True
    p_s = tf_s.paragraphs[0]
    p_s.alignment = PP_ALIGN.CENTER
    r_s = p_s.add_run()
    r_s.text = ("94% of events are unplanned  |  62% (5,030) flagged High priority  "
                "|  Event mix shifts month-over-month without a learning loop")
    r_s.font.size  = Pt(14)
    r_s.font.bold  = True
    r_s.font.color.rgb = BTP_ACCENT

    return slide


def slide_03_data(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_title_bar(slide, "The Data: Real BTP Records — No Synthetic Fabrication",
                  "8,173 ASTraM events · 6 months · entire city · used as-is from the provided dataset")
    add_accent_line(slide)

    bullets = [
        "▸  Time span: November 2023 – April 2024 — 6 full months of city-wide incident data",
        "▸  Coverage: all corridors, zones, junctions, and priority levels across Bengaluru",
        "▸  116 malformed timestamps caught and corrected during ingestion — zero silent data loss",
        "▸  Features used at inference: event cause, type, corridor, zone, junction, police station,",
        "      vehicle type, priority, latitude/longitude, hour, day-of-week, month",
        "▸  Only provided ASTraM data used — no external feeds, no synthetic rows",
        "▸  Every figure traces to a specific cell in the dataset or EDA table — fully auditable",
    ]

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.55), Inches(12.3), Inches(5.3))
    tf = tx.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(7)
        run = p.add_run()
        run.text = b
        if b.startswith("      "):
            run.font.size  = Pt(15)
            run.font.color.rgb = GREY_MID
        else:
            run.font.size  = Pt(18)
            run.font.color.rgb = BTP_TEXT

    return slide


def slide_04_insights(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_title_bar(slide, "Key Insights: What the Data Actually Says",
                  "Four findings that directly shape the solution design")
    add_accent_line(slide)

    # Left column: 4 findings
    findings = [
        ("60% breakdowns  ·  two corridors carry 100% high-priority",
         "4,896 vehicle breakdowns — concentrated on Mysore Road (743 events)\n"
         "and Bellary Road 1 (610 events), both 100% high-priority"),
        ("Construction = 5-hour bottleneck",
         "Median clearance 296 min vs 41 min for breakdowns;\n"
         "P75 stretches to 427 min — 7× slower than a typical breakdown"),
        ("80% of events are off-peak  ·  2 AM is the busiest single hour",
         "845 events at 2 AM — uncleared night breakdowns become the 8 AM gridlock.\n"
         "Night-shift coverage is not optional; it is the peak-volume hour."),
        ("Event mix is drifting — static models degrade quietly",
         "Breakdown share: 66% → 49% over 6 months.\n"
         "Road-closure rate: 6.4% → 9.4%. A frozen model never knows this."),
    ]

    tops = [Inches(1.6), Inches(2.7), Inches(3.85), Inches(5.0)]
    for i, (hdg, body) in enumerate(findings):
        lbl = slide.shapes.add_shape(1, Inches(0.35), tops[i]+Inches(0.05),
                                     Inches(0.07), Inches(0.85))
        add_solid_fill(lbl, BTP_ACCENT)
        lbl.line.fill.background()

        tx = slide.shapes.add_textbox(Inches(0.58), tops[i], Inches(7.5), Inches(1.05))
        tf = tx.text_frame; tf.word_wrap = True
        p1 = tf.paragraphs[0]
        r1 = p1.add_run(); r1.text = hdg
        r1.font.size = Pt(16); r1.font.bold = True; r1.font.color.rgb = BTP_DARK
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = body
        r2.font.size = Pt(14); r2.font.color.rgb = GREY_MID

    # Right: fig_04 corridors chart
    add_image(slide,
              os.path.join(FIG, "fig_04_corridors.png"),
              Inches(8.35), Inches(1.55), Inches(4.7), Inches(5.6))

    return slide


def slide_05_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_title_bar(slide, "Our Solution: Three Pillars + One Dashboard",
                  "One input (the event report).  Three outputs.  One interface.")
    add_accent_line(slide)

    pillars = [
        ("PILLAR 1 — Forecast Impact",
         BTP_ACCENT,
         [
             "Road-closure probability (AUC 0.816 on future data)",
             "Expected clearance time — benchmarked per cause from EDA",
             "Severity label: Low / Medium / High",
         ]),
        ("PILLAR 2 — Recommendation Engine",
         RGBColor(0x00, 0x7A, 0xCC),
         [
             "Officer count  |  Barricade flag  |  Diversion route",
             "Nearest responding police station",
             "Plain-English rationale — designed for field officers",
         ]),
        ("PILLAR 3 — Post-Event Learning Loop",
         RGBColor(0x27, 0xAE, 0x60),
         [
             "Monthly retrain on new resolved events",
             "Drift monitoring — alert if AUC drops below 0.70",
             "Feb 2024 proof: retrained +1.8 pp vs frozen static model",
         ]),
    ]

    box_tops = [Inches(1.55), Inches(3.1), Inches(4.65)]
    for i, (title, colour, items) in enumerate(pillars):
        # coloured header bar
        hdr = slide.shapes.add_shape(1,
            Inches(0.35), box_tops[i], Inches(12.6), Inches(0.45))
        add_solid_fill(hdr, colour)
        hdr.line.fill.background()
        tx_hdr = slide.shapes.add_textbox(
            Inches(0.5), box_tops[i]+Inches(0.04), Inches(12), Inches(0.38))
        tf_h = tx_hdr.text_frame
        p = tf_h.paragraphs[0]
        r = p.add_run(); r.text = title
        r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE

        # body bullets
        tx_body = slide.shapes.add_textbox(
            Inches(0.5), box_tops[i]+Inches(0.48), Inches(12.3), Inches(1.0))
        tf_b = tx_body.text_frame; tf_b.word_wrap = True
        first = True
        for it in items:
            p = tf_b.paragraphs[0] if first else tf_b.add_paragraph()
            first = False
            p.space_before = Pt(3)
            run = p.add_run(); run.text = "   ▸  " + it
            run.font.size = Pt(15); run.font.color.rgb = BTP_TEXT

    # Footer
    tx_f = slide.shapes.add_textbox(Inches(0.35), Inches(6.3), Inches(12.6), Inches(0.7))
    tf_f = tx_f.text_frame
    p_f = tf_f.paragraphs[0]; p_f.alignment = PP_ALIGN.CENTER
    r_f = p_f.add_run()
    r_f.text = ("Streamlit dashboard: live map · hotspot heatmap · trend charts · "
                "Event Simulator (live demo at the finale)")
    r_f.font.size = Pt(14); r_f.font.bold = True; r_f.font.color.rgb = BTP_ACCENT

    return slide


def slide_06_models(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_title_bar(slide, "Pillar 1: Impact Forecasting — Honest Model Results",
                  "CatBoost · time-based train/test split · trained on past, tested on future months")
    add_accent_line(slide)

    # Left: results text
    left_items = [
        "Road-closure classifier (CatBoost)  —  primary signal",
        "",
        "Future months (Mar–Apr 2024)    ROC-AUC  0.816   F1  0.450",
        "Unseen corridors (6 held out)      ROC-AUC  ~0.70 (0.696)   F1  0.335",
        "Cold-start (unseen corridor + junction)   ROC-AUC  ~0.73 (0.731)",
        "",
        "F1 is lower — road closures are rare (8.3%).  AUC is the correct metric",
        "for a skewed binary target; it measures ranking ability at all thresholds.",
        "",
        "Top feature: event cause (36 pts) — closure risk estimated from cause alone",
        "even when corridor is entirely unseen.",
        "",
        "Duration benchmark  —  honest planning uncertainty",
        "MAE ~103 min on future data.  R² ≈ 0  (duration is not reliably predictable).",
        "EDA cause medians used instead:  breakdown 41 min  ·  construction 296 min",
    ]

    tx = slide.shapes.add_textbox(Inches(0.4), Inches(1.55), Inches(7.8), Inches(5.6))
    tf = tx.text_frame; tf.word_wrap = True
    first = True
    for item in left_items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if item == "":
            p.space_before = Pt(4)
            continue
        run = p.add_run(); run.text = item
        if item.startswith("Road-closure") or item.startswith("Duration benchmark"):
            run.font.size = Pt(17); run.font.bold = True; run.font.color.rgb = BTP_DARK
        elif "0.816" in item or "0.70" in item or "0.73" in item:
            run.font.size = Pt(15); run.font.bold = True; run.font.color.rgb = BTP_ACCENT
        elif item.startswith("Top feature") or "41 min" in item:
            run.font.size = Pt(14); run.font.bold = False; run.font.color.rgb = BTP_DARK
        else:
            run.font.size = Pt(14); run.font.color.rgb = GREY_MID

    # Right: fig_07 learning curve
    add_image(slide,
              os.path.join(FIG, "fig_07_learning.png"),
              Inches(8.4), Inches(1.55), Inches(4.6), Inches(5.6))

    return slide


def slide_07_recommendation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_title_bar(slide, "Pillar 2: Recommendation Engine — A Worked Example",
                  "What a field officer sees within seconds of logging an event")
    add_accent_line(slide)

    # Scenario label
    scen_box = slide.shapes.add_shape(1, Inches(0.4), Inches(1.55), Inches(12.5), Inches(0.45))
    add_solid_fill(scen_box, BTP_DARK)
    scen_box.line.fill.background()
    tx_s = slide.shapes.add_textbox(Inches(0.55), Inches(1.58), Inches(12), Inches(0.38))
    tf_s = tx_s.text_frame
    p_s = tf_s.paragraphs[0]
    r_s = p_s.add_run()
    r_s.text = "Scenario: Construction on Mysore Road · 9 AM · High priority"
    r_s.font.size = Pt(16); r_s.font.bold = True; r_s.font.color.rgb = WHITE

    # Output fields
    fields = [
        ("Expected clearance",     "~296 min  (EDA: construction median — nearly 5 hours)"),
        ("Road-closure probability","62%"),
        ("Severity",               "HIGH  (all four scoring signals fire)"),
        ("Recommended officers",   "10  (base 6 + 2 peak-hour + 2 closure bonus)"),
        ("Barricading",            "YES  (62% > 35% threshold = 4× the historical base rate)"),
        ("Diversion",              "Magadi Road or Chord Road"),
        ("Responding station",     "Halasuru Gate"),
    ]

    tops = [Inches(2.12 + i * 0.62) for i in range(len(fields))]
    for i, (lbl, val) in enumerate(fields):
        # label
        tx_l = slide.shapes.add_textbox(Inches(0.5), tops[i], Inches(4.0), Inches(0.55))
        tf_l = tx_l.text_frame
        p_l = tf_l.paragraphs[0]
        r_l = p_l.add_run(); r_l.text = lbl
        r_l.font.size = Pt(16); r_l.font.bold = True; r_l.font.color.rgb = BTP_DARK
        # value
        tx_v = slide.shapes.add_textbox(Inches(4.6), tops[i], Inches(8.5), Inches(0.55))
        tf_v = tx_v.text_frame
        p_v = tf_v.paragraphs[0]
        r_v = p_v.add_run(); r_v.text = val
        r_v.font.size = Pt(16); r_v.font.color.rgb = BTP_ACCENT

    # Logic note at bottom
    note_items = [
        "Officer count: severity tier sets base (Low=2, Med=4, High=6); peak-hour and closure bonuses are structured — no black box",
        "Barricade threshold 35% = 4× historical base rate of 8.3%",
        "Station fallback: corridor history → zone history → Haversine nearest — never crashes on unseen locations",
    ]
    tx_n = slide.shapes.add_textbox(Inches(0.4), Inches(6.3), Inches(12.6), Inches(0.9))
    tf_n = tx_n.text_frame; tf_n.word_wrap = True
    first = True
    for note in note_items:
        p = tf_n.paragraphs[0] if first else tf_n.add_paragraph()
        first = False
        run = p.add_run(); run.text = "▸  " + note
        run.font.size = Pt(12); run.font.color.rgb = GREY_MID

    return slide


def slide_08_learning(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_title_bar(slide, "Pillar 3: The Learning Loop",
                  "A static model is a liability — the event mix is already shifting")
    add_accent_line(slide)

    # Table header
    col_heads = ["Month tested", "Static (frozen) AUC", "Retrained AUC", "Difference"]
    col_x = [Inches(0.4), Inches(3.2), Inches(7.0), Inches(10.2)]
    col_w = [Inches(2.7), Inches(3.6), Inches(3.0), Inches(2.8)]

    hdr_top = Inches(1.62)
    for j, (h, cx, cw) in enumerate(zip(col_heads, col_x, col_w)):
        hdr_cell = slide.shapes.add_shape(1, cx, hdr_top, cw, Inches(0.38))
        add_solid_fill(hdr_cell, BTP_DARK)
        hdr_cell.line.fill.background()
        tx = slide.shapes.add_textbox(cx+Inches(0.05), hdr_top+Inches(0.03),
                                      cw-Inches(0.1), Inches(0.32))
        tf = tx.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h
        r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE

    # Table rows
    rows = [
        ("Jan 2024", "0.751", "0.747", "−0.004"),
        ("Feb 2024", "0.696", "0.714", "+0.018  ← key failure"),
        ("Mar 2024", "0.792", "0.805", "+0.013"),
        ("Apr 2024", "0.822", "0.838", "+0.016"),
    ]
    highlight_row = 1  # Feb

    for i, (mo, static, retrain, delta) in enumerate(rows):
        row_top = Inches(2.05 + i * 0.52)
        vals = [mo, static, retrain, delta]
        bg_col = RGBColor(0xFF, 0xEF, 0xE0) if i == highlight_row else RGBColor(0xF0, 0xF3, 0xF7)
        for j, (v, cx, cw) in enumerate(zip(vals, col_x, col_w)):
            cell = slide.shapes.add_shape(1, cx, row_top, cw, Inches(0.47))
            add_solid_fill(cell, bg_col)
            cell.line.fill.background()
            tx = slide.shapes.add_textbox(cx+Inches(0.05), row_top+Inches(0.04),
                                          cw-Inches(0.1), Inches(0.38))
            tf = tx.text_frame
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = v
            r.font.size = Pt(15)
            if i == highlight_row and j > 0:
                r.font.bold = True
                r.font.color.rgb = BTP_ACCENT
            else:
                r.font.color.rgb = BTP_TEXT

    # Right: fig_07
    add_image(slide,
              os.path.join(FIG, "fig_07_learning.png"),
              Inches(8.1), Inches(1.62), Inches(4.9), Inches(4.2))

    # Key takeaway bullets
    bullets = [
        "Feb is the clearest case: static drops to AUC 0.696; retrained holds at 0.714",
        "Retraining never materially hurts (Jan: −0.004 pp, negligible)",
        "Recommended cadence: monthly — aligns with observed cause-mix shift",
        "Alert triggered automatically if AUC drops below 0.70",
    ]
    tx_b = slide.shapes.add_textbox(Inches(0.4), Inches(5.05), Inches(7.5), Inches(2.1))
    tf_b = tx_b.text_frame; tf_b.word_wrap = True
    first = True
    for b in bullets:
        p = tf_b.paragraphs[0] if first else tf_b.add_paragraph()
        first = False
        p.space_before = Pt(5)
        run = p.add_run(); run.text = "▸  " + b
        run.font.size = Pt(15); run.font.color.rgb = BTP_TEXT

    return slide


def slide_09_robustness(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_title_bar(slide, "Robustness: Works on Data the Model Has Never Seen",
                  "Three validation dimensions — because production is always the future")
    add_accent_line(slide)

    dimensions = [
        ("Future time (primary)",
         "Train Nov–Feb  →  test Mar–Apr 2024",
         "AUC 0.816 (closure)  ·  MAE ~103 min (duration)  ·  R² ≈ 0 (intentional — use EDA medians)"),
        ("Unseen locations",
         "6 corridors held out entirely from training",
         "AUC ~0.70 (0.696)  ·  Cold-start (unseen corridor + junction): AUC ~0.73 (0.731)"),
        ("Completely new event",
         "Fog / Low Visibility on NH-44 Bypass — not in any training record",
         "Degrades gracefully: zone fallback station  ·  AUC-estimated closure  ·  41-min benchmark"),
    ]

    tops = [Inches(1.62), Inches(2.9), Inches(4.18)]
    for i, (dim, test, result) in enumerate(dimensions):
        lbl = slide.shapes.add_shape(1, Inches(0.35), tops[i]+Inches(0.08),
                                     Inches(0.07), Inches(0.95))
        add_solid_fill(lbl, BTP_ACCENT)
        lbl.line.fill.background()

        tx = slide.shapes.add_textbox(Inches(0.6), tops[i], Inches(12.3), Inches(1.15))
        tf = tx.text_frame; tf.word_wrap = True
        p1 = tf.paragraphs[0]
        r1 = p1.add_run(); r1.text = dim
        r1.font.size = Pt(17); r1.font.bold = True; r1.font.color.rgb = BTP_DARK
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = test
        r2.font.size = Pt(14); r2.font.italic = True; r2.font.color.rgb = GREY_MID
        p3 = tf.add_paragraph()
        r3 = p3.add_run(); r3.text = result
        r3.font.size = Pt(15); r3.font.color.rgb = BTP_ACCENT

    # Built-in robustness bullets
    built_in = [
        "Out-of-fold target encoding (alpha=20) prevents leakage; new corridors fall back to global prior — no crash",
        "CatBoost handles completely unseen cause labels natively at inference time",
        "Three-layer station lookup (corridor → zone → Haversine) always returns a valid station",
        "116 malformed timestamps caught and corrected — pipeline tolerates real-world messy ASTraM data",
        "78/78 formal tests pass — including empty dict, all-unseen categories, and sparse inputs",
    ]
    tx_b = slide.shapes.add_textbox(Inches(0.4), Inches(5.4), Inches(12.6), Inches(1.85))
    tf_b = tx_b.text_frame; tf_b.word_wrap = True
    first = True
    for b in built_in:
        p = tf_b.paragraphs[0] if first else tf_b.add_paragraph()
        first = False
        p.space_before = Pt(4)
        run = p.add_run(); run.text = "▸  " + b
        run.font.size = Pt(14); run.font.color.rgb = BTP_TEXT

    return slide


def slide_10_impact(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_title_bar(slide, "Real-World Impact & Scalability for BTP",
                  "What changes the day BTP deploys this system")
    add_accent_line(slide)

    left_col = [
        ("Immediate operational changes", [
            "Every logged event gets severity label + officer count in under 1 second",
            "Closure-risk events (Mysore Road construction 62%) trigger barricade pre-staging",
            "Night-shift gap quantified: 845 events at 2 AM demand staffed response",
            "Seasonal planning: March surge (1,956 events) can be resourced in advance",
        ]),
        ("Scalability — honest assessment", [
            "Retraining: under 5 minutes on a laptop; new corridors handled without code changes",
            "Real-time feed: replace CSV loader with API call — all downstream logic unchanged",
            "Other cities: retrain on that city's data; one config change (fallback station)",
        ]),
    ]

    right_col = [
        ("What this is not", [
            "Not a replacement for field judgment — officer counts are starting points",
            "Not a real-time signal controller — it advises dispatchers, does not automate",
            "Not a substitute for live traffic feed data — diversion routes not dynamically verified",
        ]),
    ]

    top = Inches(1.6)
    for heading, items in left_col:
        hdr = slide.shapes.add_textbox(Inches(0.4), top, Inches(7.0), Inches(0.42))
        tf_h = hdr.text_frame
        p_h = tf_h.paragraphs[0]
        r_h = p_h.add_run(); r_h.text = heading
        r_h.font.size = Pt(17); r_h.font.bold = True; r_h.font.color.rgb = BTP_DARK
        top += Inches(0.45)
        for item in items:
            tx_i = slide.shapes.add_textbox(Inches(0.5), top, Inches(6.8), Inches(0.42))
            tf_i = tx_i.text_frame; tf_i.word_wrap = True
            p_i = tf_i.paragraphs[0]
            r_i = p_i.add_run(); r_i.text = "▸  " + item
            r_i.font.size = Pt(14); r_i.font.color.rgb = BTP_TEXT
            top += Inches(0.44)
        top += Inches(0.15)

    right_top = Inches(1.6)
    for heading, items in right_col:
        hdr = slide.shapes.add_textbox(Inches(7.8), right_top, Inches(5.2), Inches(0.42))
        tf_h = hdr.text_frame
        p_h = tf_h.paragraphs[0]
        r_h = p_h.add_run(); r_h.text = heading
        r_h.font.size = Pt(17); r_h.font.bold = True; r_h.font.color.rgb = RGBColor(0xC0, 0x39, 0x2B)
        right_top += Inches(0.45)
        for item in items:
            tx_i = slide.shapes.add_textbox(Inches(7.9), right_top, Inches(5.0), Inches(0.5))
            tf_i = tx_i.text_frame; tf_i.word_wrap = True
            p_i = tf_i.paragraphs[0]
            r_i = p_i.add_run(); r_i.text = "✗  " + item
            r_i.font.size = Pt(14); r_i.font.color.rgb = GREY_MID
            right_top += Inches(0.52)

    return slide


def slide_11_tests(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_title_bar(slide, "Test Suite: 78/78 Passing — Zero Failures",
                  "Formal reproducible test suite covers every layer of the system")
    add_accent_line(slide)

    sections = [
        ("Data & Cleaning edge cases",  "32 tests",
         "Timestamp parsing · negative durations · unseen cause labels · missing fields"),
        ("Accuracy on unseen data",     "7 tests",
         "AUC on future months · unseen corridors · cold-start — no training data leakage"),
        ("predict_impact robustness",   "12 tests",
         "Empty dict · all-unseen categories · sparse input — always returns valid in-range output"),
        ("recommend() robustness",      "24 tests",
         "All 17 cause values · unseen corridor + cause · empty dict — always returns all 8 keys"),
        ("Dashboard smoke test",        "3 tests",
         "AppTest load · form submit · output rendered — zero exceptions on a full run"),
    ]

    tops = [Inches(1.62 + i * 0.96) for i in range(len(sections))]
    for i, (section, count, detail) in enumerate(sections):
        bg = slide.shapes.add_shape(1, Inches(0.35), tops[i], Inches(12.6), Inches(0.88))
        col = RGBColor(0xF0, 0xF8, 0xF1) if i % 2 == 0 else RGBColor(0xF5, 0xF7, 0xFA)
        add_solid_fill(bg, col)
        bg.line.fill.background()

        tx = slide.shapes.add_textbox(Inches(0.55), tops[i]+Inches(0.04), Inches(11.5), Inches(0.8))
        tf = tx.text_frame; tf.word_wrap = True
        p1 = tf.paragraphs[0]
        r1 = p1.add_run(); r1.text = f"({count})  {section}"
        r1.font.size = Pt(17); r1.font.bold = True; r1.font.color.rgb = BTP_DARK
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = detail
        r2.font.size = Pt(14); r2.font.color.rgb = GREY_MID

        # green tick
        tx_tick = slide.shapes.add_textbox(Inches(12.3), tops[i]+Inches(0.15),
                                           Inches(0.6), Inches(0.55))
        tf_t = tx_tick.text_frame
        p_t = tf_t.paragraphs[0]; p_t.alignment = PP_ALIGN.CENTER
        r_t = p_t.add_run(); r_t.text = "✓"
        r_t.font.size = Pt(24); r_t.font.bold = True
        r_t.font.color.rgb = RGBColor(0x27, 0xAE, 0x60)

    return slide


def slide_12_ask(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide)
    add_title_bar(slide, "The Ask",
                  "One prototype ready now.  Three things needed to make it operational.")
    add_accent_line(slide)

    # What we've built
    built = slide.shapes.add_textbox(Inches(0.4), Inches(1.6), Inches(12.5), Inches(0.38))
    tf_b = built.text_frame
    p_b = tf_b.paragraphs[0]
    r_b = p_b.add_run(); r_b.text = "What we've built"
    r_b.font.size = Pt(18); r_b.font.bold = True; r_b.font.color.rgb = BTP_DARK

    built_items = [
        "Notebook: full EDA + model training + learning-loop comparison, all reproducible",
        "Streamlit dashboard with live Event Simulator — inference under 1 second per event",
        "Monthly retraining script ready to schedule; 78/78 formal tests pass",
    ]
    tx_bi = slide.shapes.add_textbox(Inches(0.5), Inches(2.02), Inches(12.3), Inches(1.0))
    tf_bi = tx_bi.text_frame; tf_bi.word_wrap = True
    first = True
    for it in built_items:
        p = tf_bi.paragraphs[0] if first else tf_bi.add_paragraph()
        first = False
        p.space_before = Pt(4)
        run = p.add_run(); run.text = "▸  " + it
        run.font.size = Pt(15); run.font.color.rgb = BTP_TEXT

    # BTP asks
    btp_items = [
        ("1", "Access to a live ASTraM feed",
         "Move from historical replay to real-time recommendations — pipeline is already built for it"),
        ("2", "Feedback labels from field officers",
         '"Was the closure prediction right? Was manpower sufficient?" — closes the learning loop'),
        ("3", "Pilot corridor: Mysore Road  ·  30 days",
         "Run recommendations alongside manual dispatch; compare officer utilisation and response time"),
    ]

    tops_ask = [Inches(3.25), Inches(4.35), Inches(5.45)]
    for i, (num, heading, detail) in enumerate(btp_items):
        num_box = slide.shapes.add_shape(1, Inches(0.4), tops_ask[i], Inches(0.5), Inches(0.5))
        add_solid_fill(num_box, BTP_ACCENT)
        num_box.line.fill.background()
        tx_n = slide.shapes.add_textbox(Inches(0.4), tops_ask[i]+Inches(0.04),
                                        Inches(0.5), Inches(0.42))
        tf_n = tx_n.text_frame
        p_n = tf_n.paragraphs[0]; p_n.alignment = PP_ALIGN.CENTER
        r_n = p_n.add_run(); r_n.text = num
        r_n.font.size = Pt(16); r_n.font.bold = True; r_n.font.color.rgb = WHITE

        tx_a = slide.shapes.add_textbox(Inches(1.05), tops_ask[i], Inches(11.6), Inches(0.9))
        tf_a = tx_a.text_frame; tf_a.word_wrap = True
        p1 = tf_a.paragraphs[0]
        r1 = p1.add_run(); r1.text = heading
        r1.font.size = Pt(17); r1.font.bold = True; r1.font.color.rgb = BTP_DARK
        p2 = tf_a.add_paragraph()
        r2 = p2.add_run(); r2.text = detail
        r2.font.size = Pt(14); r2.font.color.rgb = GREY_MID

    # Flipkart ask
    fk = slide.shapes.add_textbox(Inches(0.4), Inches(6.55), Inches(12.6), Inches(0.62))
    tf_fk = fk.text_frame; tf_fk.word_wrap = True
    p_fk = tf_fk.paragraphs[0]; p_fk.alignment = PP_ALIGN.CENTER
    r_fk = p_fk.add_run()
    r_fk.text = ("Flipkart ask: infrastructure + deployment support for the Streamlit dashboard "
                 "on BTP's internal network")
    r_fk.font.size = Pt(14); r_fk.font.bold = True; r_fk.font.color.rgb = BTP_ACCENT

    return slide


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_data(prs)
    slide_04_insights(prs)
    slide_05_solution(prs)
    slide_06_models(prs)
    slide_07_recommendation(prs)
    slide_08_learning(prs)
    slide_09_robustness(prs)
    slide_10_impact(prs)
    slide_11_tests(prs)
    slide_12_ask(prs)

    out = os.path.join(BASE, "presentation.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slide count: {len(prs.slides)}")
    return out


if __name__ == "__main__":
    build()
